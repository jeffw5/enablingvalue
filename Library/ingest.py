#!/usr/bin/env python3
"""
ingest.py — Knowledge Base Ingestion Pipeline
The Value Enablement Group, LLC

Supports: .md, .txt, .pdf, .docx, .pptx, .html, .json

Usage:
    python3 ingest.py

Requirements:
    pip3 install pinecone requests python-dotenv pypdf \
                 python-docx python-pptx beautifulsoup4 \
                 markdown --break-system-packages

Setup:
    Create a .env file with:
        PINECONE_API_KEY=pcsk_7XhdCs_CmYjSLBVx9PthfdboZveK8MCAAjK6Tr7fxmrvgPg99od8cYZsw1p6heJPPAGSYr
        PINECONE_HOST=https://enabling-value-xxx.svc.aped-4627-b74a.pinecone.io
        GRAPHDB_ENDPOINT=https://your-tunnel.trycloudflare.com/repositories/Value-kb
"""

import os
import re
import json
import hashlib
import datetime
import requests
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_HOST    = os.getenv("PINECONE_HOST")
GRAPHDB_ENDPOINT = os.getenv("GRAPHDB_ENDPOINT")
KNOWLEDGE_BASE   = Path("knowledge-base")

PC_HEADERS = {
    "Api-Key": PINECONE_API_KEY,
    "Content-Type": "application/json",
    "X-Pinecone-API-Version": "2024-07"
}

CHUNK_SIZE    = 500
CHUNK_OVERLAP = 50
BASE_URI      = "https://enablingvalue.com/kb#"

SUPPORTED = {".md", ".txt", ".pdf", ".docx", ".pptx", ".html", ".json"}


# ──────────────────────────────────────────
# TEXT EXTRACTION — one function per format
# ──────────────────────────────────────────

def extract_md(filepath):
    return filepath.read_text(encoding="utf-8", errors="ignore")

def extract_txt(filepath):
    return filepath.read_text(encoding="utf-8", errors="ignore")

def extract_pdf(filepath):
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(filepath))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n".join(pages)
    except ImportError:
        print("    ⚠ pypdf not installed: pip3 install pypdf --break-system-packages")
        return ""
    except Exception as e:
        print(f"    ⚠ PDF error: {e}")
        return ""

def extract_docx(filepath):
    try:
        from docx import Document
        doc = Document(str(filepath))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs)
    except ImportError:
        print("    ⚠ python-docx not installed: pip3 install python-docx --break-system-packages")
        return ""
    except Exception as e:
        print(f"    ⚠ DOCX error: {e}")
        return ""

def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        print("    ⚠ python-pptx not installed: pip3 install python-pptx --break-system-packages")
        return ""
    except Exception as e:
        print(f"    ⚠ PPTX error: {e}")
        return ""

def extract_html(filepath):
    try:
        from bs4 import BeautifulSoup
        html = filepath.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        # Remove script and style tags
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        print("    ⚠ beautifulsoup4 not installed: pip3 install beautifulsoup4 --break-system-packages")
        return ""
    except Exception as e:
        print(f"    ⚠ HTML error: {e}")
        return ""

def extract_json(filepath):
    try:
        data = json.loads(filepath.read_text(encoding="utf-8"))
        # Flatten JSON to readable text
        return json.dumps(data, indent=2)
    except Exception as e:
        print(f"    ⚠ JSON error: {e}")
        return ""

EXTRACTORS = {
    ".md":   extract_md,
    ".txt":  extract_txt,
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".json": extract_json,
}

def extract_text(filepath: Path) -> str:
    suffix = filepath.suffix.lower()
    extractor = EXTRACTORS.get(suffix)
    if extractor:
        return extractor(filepath)
    print(f"    ⚠ No extractor for {suffix}")
    return ""


# ──────────────────────────────────────────
# CHUNKING
# ──────────────────────────────────────────

def chunk_text(text: str) -> list:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i : i + CHUNK_SIZE])
        if chunk.strip():
            chunks.append(chunk.strip())
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


# ──────────────────────────────────────────
# METADATA
# ──────────────────────────────────────────

def extract_frontmatter(text: str):
    meta = {}
    if text.startswith("---"):
        end = text.find("---", 3)
        if end != -1:
            fm = text[3:end].strip()
            body = text[end+3:].strip()
            for line in fm.splitlines():
                if ":" in line:
                    k, _, v = line.partition(":")
                    meta[k.strip()] = v.strip()
            return meta, body
    return meta, text

def infer_metadata(filepath: Path) -> dict:
    parts = filepath.parts
    domain = parts[1] if len(parts) > 1 else "general"
    return {
        "domain":    domain,
        "filename":  filepath.name,
        "filepath":  str(filepath),
        "filetype":  filepath.suffix.lower().lstrip("."),
        "source":    "knowledge-base",
        "graph_uri": f"{BASE_URI}{domain}/{filepath.stem}"
    }


# ──────────────────────────────────────────
# PINECONE
# ──────────────────────────────────────────

def upsert_to_pinecone(vectors: list) -> bool:
    if not vectors:
        return True

    # Step 1 — Get embeddings from Pinecone inference API
    embed_payload = {
        "inputs": [{"text": v["text"]} for v in vectors],
        "model": "llama-text-embed-v2"
    }

    try:
        r = requests.post(
            "https://api.pinecone.io/embed",
            headers={
                "Api-Key": PINECONE_API_KEY,
                "Content-Type": "application/json",
                "X-Pinecone-API-Version": "2024-07"
            },
            json=embed_payload,
            timeout=60
        )

        if r.status_code != 200:
            print(f"    ⚠ Embed error {r.status_code}: {r.text[:300]}")
            return False

        embeddings = r.json().get("data", [])
        if not embeddings:
            print("    ⚠ No embeddings returned")
            return False

        # Step 2 — Build upsert records
        records = []
        for v, emb in zip(vectors, embeddings):
            records.append({
                "id":       v["id"],
                "values":   emb["values"],
                "metadata": v["metadata"]
            })

        # Step 3 — Upsert to index
        r2 = requests.post(
            f"{PINECONE_HOST}/vectors/upsert",
            headers=PC_HEADERS,
            json={"vectors": records},
            timeout=60
        )

        if r2.status_code == 200:
            return True
        else:
            print(f"    ⚠ Upsert error {r2.status_code}: {r2.text[:300]}")
            return False

    except Exception as e:
        print(f"    ⚠ Pinecone error: {e}")
        return False


# ──────────────────────────────────────────
# GRAPHDB
# ──────────────────────────────────────────

def log_to_graphdb(meta: dict, chunk_count: int) -> bool:
    if not GRAPHDB_ENDPOINT:
        return True

    uri       = meta["graph_uri"]
    domain    = meta["domain"]
    filename  = meta["filename"].replace('"', '\\"')
    filepath  = meta["filepath"].replace('"', '\\"')
    filetype  = meta["filetype"]
    timestamp = datetime.datetime.utcnow().isoformat()

    sparql = f"""
PREFIX ev: <{BASE_URI}>
PREFIX dc: <http://purl.org/dc/elements/1.1/>
PREFIX xsd: <http://www.w3.org/2001/XMLSchema#>

INSERT DATA {{
    <{uri}> a ev:KnowledgeDocument ;
        dc:title "{filename}" ;
        ev:domain "{domain}" ;
        ev:filetype "{filetype}" ;
        ev:chunkCount "{chunk_count}"^^xsd:integer ;
        ev:filepath "{filepath}" ;
        ev:ingestedAt "{timestamp}"^^xsd:dateTime .
}}
"""

    try:
        r = requests.post(
            f"{GRAPHDB_ENDPOINT}/statements",
            data=sparql.encode("utf-8"),
            headers={"Content-Type": "application/sparql-update"},
            timeout=30
        )
        return r.status_code in (200, 204)
    except Exception as e:
        print(f"    ⚠ GraphDB error: {e}")
        return False


# ──────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────

def generate_id(filepath: str, chunk_index: int) -> str:
    return hashlib.md5(f"{filepath}::{chunk_index}".encode()).hexdigest()


def process_file(filepath: Path) -> int:
    print(f"  [{filepath.suffix.upper().lstrip('.')}] {filepath}")

    raw_text = extract_text(filepath)
    if not raw_text.strip():
        print(f"    ⚠ No text extracted — skipping")
        return 0

    fm_meta, body = extract_frontmatter(raw_text)
    meta = infer_metadata(filepath)
    meta.update(fm_meta)

    chunks = chunk_text(body)
    if not chunks:
        return 0

    print(f"    → {len(chunks)} chunks")

    batch_size = 20
    total = 0

    for i in range(0, len(chunks), batch_size):
        batch = chunks[i : i + batch_size]
        vectors = []
        for j, chunk in enumerate(batch):
            idx = i + j
            vectors.append({
                "id":   generate_id(str(filepath), idx),
                "text": chunk,
                "metadata": {
                    "text":      chunk[:500],
                    "domain":    meta.get("domain", "general"),
                    "filename":  meta.get("filename", ""),
                    "filepath":  meta.get("filepath", ""),
                    "filetype":  meta.get("filetype", ""),
                    "graph_uri": meta.get("graph_uri", ""),
                    "chunk_idx": idx,
                    "tags":      str(meta.get("tags", ""))
                }
            })

        ok = upsert_to_pinecone(vectors)
        if ok:
            total += len(batch)
            print(f"    ✓ Batch {i//batch_size + 1} → {len(batch)} chunks")
        else:
            print(f"    ✗ Batch {i//batch_size + 1} failed")

    if GRAPHDB_ENDPOINT:
        ok = log_to_graphdb(meta, total)
        print(f"    {'✓' if ok else '⚠'} GraphDB {'logged' if ok else 'log failed'}")

    return total


def main():
    print("=" * 60)
    print("  VEG Knowledge Base Ingestion Pipeline")
    print("=" * 60)

    if not PINECONE_API_KEY:
        print("❌ PINECONE_API_KEY missing from .env"); return
    if not PINECONE_HOST:
        print("❌ PINECONE_HOST missing from .env"); return

    print(f"✓ Pinecone: {PINECONE_HOST}")
    print(f"✓ GraphDB:  {GRAPHDB_ENDPOINT or 'not configured'}")
    print(f"✓ Source:   {KNOWLEDGE_BASE.absolute()}")
    print()

    if not KNOWLEDGE_BASE.exists():
        print(f"❌ knowledge-base/ not found at {KNOWLEDGE_BASE.absolute()}")
        print("  Run this script from your enablingvalue repo root")
        return

    files = sorted([
        f for f in KNOWLEDGE_BASE.rglob("*")
        if f.is_file()
        and f.suffix.lower() in SUPPORTED
        and not f.name.startswith(".")
        and f.name.lower() not in ("readme.md", "readme.txt")
    ])

    if not files:
        print("⚠ No content files found in knowledge-base/")
        return

    # Show what we found
    by_type = {}
    for f in files:
        t = f.suffix.lower()
        by_type[t] = by_type.get(t, 0) + 1
    print(f"Found {len(files)} files:")
    for t, n in sorted(by_type.items()):
        print(f"  {t}: {n}")
    print()

    total_chunks = 0
    errors = []

    for filepath in files:
        try:
            n = process_file(filepath)
            total_chunks += n
            print()
        except Exception as e:
            errors.append((filepath, str(e)))
            print(f"    ✗ Unexpected error: {e}\n")

    print("=" * 60)
    print(f"  ✓ {len(files)} files → {total_chunks} total chunks ingested")
    if errors:
        print(f"  ✗ {len(errors)} files failed:")
        for f, e in errors:
            print(f"    {f}: {e}")
    print("=" * 60)


if __name__ == "__main__":
    main()
